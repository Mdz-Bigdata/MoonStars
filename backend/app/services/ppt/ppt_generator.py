"""
PPT 生成器
使用 python-pptx 生成 PPTX 文件。
支持标准模式（结构化生成）和混合模式（AI背景+文字叠加）。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import List, Optional
import io
import logging
from .models import PPTOutline, Section, ImageInfo

logger = logging.getLogger(__name__)


class PPTGeneratorService:
    """
    PPT 生成服务

    支持两种生成方式：
    1. generate() - 标准结构化生成（文字+图表+配图）
    2. generate_from_images() - 全图式生成（AI 生成的完整幻灯片图片作为页面）
    3. generate_hybrid() - 混合模式（AI 背景 + 文字叠加）
    """

    def __init__(self, theme: str = "light"):
        self.prs = Presentation()
        # 设置 16:9 宽高比
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        self.theme = theme

        # 主题 Tokens - 丰富的主题选项
        self.themes = {
            # 经典主题
            "light": {"bg": "FFFFFF", "text": "1E293B", "primary": "2563EB", "accent": "3B82F6", "secondary": "F1F5F9"},
            "dark": {"bg": "0F172A", "text": "F8FAFC", "primary": "38BDF8", "accent": "0EA5E9", "secondary": "1E293B"},
            "corporate": {"bg": "FFFFFF", "text": "111827", "primary": "1E3A8A", "accent": "FBBF24", "secondary": "F3F4F6"},
            # 现代渐变风格
            "gradient": {"bg": "FAFBFF", "text": "1F2937", "primary": "6366F1", "accent": "EC4899", "secondary": "EEF2FF"},
            # 极简风格
            "minimalist": {"bg": "FFFFFF", "text": "18181B", "primary": "18181B", "accent": "A1A1AA", "secondary": "FAFAFA"},
            # 科技风格
            "tech": {"bg": "030712", "text": "E5E7EB", "primary": "22D3EE", "accent": "A855F7", "secondary": "111827"},
            # 创意多彩
            "creative": {"bg": "FFF7ED", "text": "1C1917", "primary": "F97316", "accent": "10B981", "secondary": "FFEDD5"},
            # 自然清新
            "nature": {"bg": "F0FDF4", "text": "14532D", "primary": "16A34A", "accent": "84CC16", "secondary": "DCFCE7"},
            # 商务蓝
            "business": {"bg": "F8FAFC", "text": "0F172A", "primary": "0369A1", "accent": "0891B2", "secondary": "E0F2FE"},
            # 优雅紫
            "elegant": {"bg": "FAF5FF", "text": "3B0764", "primary": "7C3AED", "accent": "C084FC", "secondary": "EDE9FE"},
        }
        self.colors = self.themes.get(theme, self.themes["light"])

    # ==================== 标准模式（保持原有逻辑） ====================

    def generate(
        self,
        outline: PPTOutline,
        slide_images: List[Optional[bytes]] = [],
        mermaid_images: List[Optional[bytes]] = [],
        fragment_images: List[Optional[bytes]] = [],
    ) -> io.BytesIO:
        """标准结构化 PPT 生成（原有逻辑完全不变）"""
        # Title Slide
        self._add_title_slide(outline.title, outline.subtitle)

        # Content Slides
        for i, section in enumerate(outline.sections):
            self._add_content_slide(
                section,
                image_data=slide_images[i] if i < len(slide_images) else None,
                mermaid_image=mermaid_images[i] if i < len(mermaid_images) else None,
                fragment_image=fragment_images[i] if i < len(fragment_images) else None,
            )

        if outline.conclusion:
            self._add_conclusion_slide(outline.conclusion)

        self._add_thank_you_slide()

        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output

    # ==================== 全图模式（banana-slides AI 生成） ====================

    def generate_from_images(self, images: List[bytes], title: str = "") -> io.BytesIO:
        """
        全图式 PPT 生成。
        将 AI 生成的完整幻灯片图片直接作为页面背景。

        Args:
            images: 图片字节数据列表（每页一张）
            title: PPT 标题（用于文件属性）

        Returns:
            PPTX BytesIO
        """
        for img_data in images:
            if img_data is None:
                continue
            try:
                slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局
                img_stream = io.BytesIO(img_data)
                slide.shapes.add_picture(
                    img_stream,
                    Inches(0),
                    Inches(0),
                    width=self.prs.slide_width,
                    height=self.prs.slide_height,
                )
            except Exception as e:
                logger.warning(f"全图模式添加页面失败: {e}")

        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output

    # ==================== 混合模式（AI 背景 + 结构化文字叠加） ====================

    def generate_hybrid(
        self,
        outline: PPTOutline,
        background_images: List[Optional[bytes]] = [],
        mermaid_images: List[Optional[bytes]] = [],
    ) -> io.BytesIO:
        """
        混合模式 PPT 生成。
        AI 背景图 + python-pptx 结构化文字叠加，兼顾美观和可编辑性。

        Args:
            outline: PPT 大纲
            background_images: 每页的背景图
            mermaid_images: Mermaid 图表图片

        Returns:
            PPTX BytesIO
        """
        # Title Slide（带背景）
        bg_img = background_images[0] if background_images else None
        self._add_hybrid_title_slide(outline.title, outline.subtitle, bg_img)

        # Content Slides
        for i, section in enumerate(outline.sections):
            bg = background_images[i + 1] if (i + 1) < len(background_images) else None
            m_img = mermaid_images[i] if i < len(mermaid_images) else None
            self._add_hybrid_content_slide(section, background_image=bg, mermaid_image=m_img)

        if outline.conclusion:
            self._add_conclusion_slide(outline.conclusion)

        self._add_thank_you_slide()

        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output

    # ==================== 私有方法 - 标准模式 ====================

    def _add_title_slide(self, title: str, subtitle: Optional[str]):
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(self.colors["bg"])

        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.colors["text"])
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True

        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.colors["primary"])
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)

    def _add_content_slide(
        self,
        section: Section,
        image_data: Optional[bytes] = None,
        mermaid_image: Optional[bytes] = None,
        fragment_image: Optional[bytes] = None,
    ):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Title Only

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(self.colors["bg"])

        title_shape = slide.shapes.title
        title_shape.text = section.title
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.colors["text"])
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True

        has_visual = image_data or mermaid_image or fragment_image
        num_points = len(section.points)
        has_table = hasattr(section, "table") and section.table

        is_visual_heavy = has_visual and (
            num_points <= 2 or section.visual_type in ["mindmap", "architecture", "flowchart"]
        )

        # 布局计算
        left = Inches(0.5)
        top = Inches(1.3)
        width = Inches(5.8) if has_visual else Inches(12.0)
        height = Inches(5.8)

        if has_table and not fragment_image:
            self._render_table(slide, section.table, left, top, width, height)
        else:
            if is_visual_heavy:
                top = Inches(6.2)
                height = Inches(1.0)
                width = Inches(12.3)

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True

            font_size = Pt(18)
            if is_visual_heavy:
                font_size = Pt(14)
            elif num_points > 8:
                font_size = Pt(14)

            for point in section.points:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                p.font.size = font_size
                p.font.color.rgb = RGBColor.from_string(self.colors["text"])

        # 视觉资源
        if has_visual:
            visual_data = fragment_image if fragment_image else (mermaid_image if mermaid_image else image_data)
            try:
                image_stream = io.BytesIO(visual_data)
                if is_visual_heavy:
                    slide.shapes.add_picture(image_stream, Inches(1.0), Inches(1.2), width=Inches(11.3))
                else:
                    slide.shapes.add_picture(image_stream, Inches(6.5), Inches(1.5), width=Inches(6.3))
            except Exception as e:
                logger.warning(f"Error inserting image: {e}")

        # 页码
        slide_idx = list(self.prs.slides).index(slide)
        if slide_idx > 0:
            footer_box = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4))
            p = footer_box.text_frame.paragraphs[0]
            p.text = str(slide_idx + 1)
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor.from_string(self.colors["text"])
            p.alignment = PP_ALIGN.RIGHT

    def _render_table(self, slide, table_md: str, left, top, width, height):
        """解析 Markdown 表格并渲染到 PPT"""
        try:
            lines = [l.strip() for l in table_md.strip().split("\n") if "|" in l]
            if not lines:
                return

            content_lines = [l for l in lines if not all(c in "|- " for c in l.replace("|", ""))]

            rows_data = []
            for line in content_lines:
                cells = [c.strip() for c in line.split("|") if c.strip() or line.count("|") > 1]
                if cells:
                    rows_data.append(cells)

            if not rows_data:
                return

            rows = len(rows_data)
            cols = max(len(r) for r in rows_data)

            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            for r_idx, row in enumerate(rows_data):
                for c_idx, cell_text in enumerate(row):
                    if c_idx < cols:
                        cell = table.cell(r_idx, c_idx)
                        cell.text = cell_text
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(12)
        except Exception as e:
            logger.warning(f"Error rendering table: {e}")

    def _add_conclusion_slide(self, conclusion: List[str]):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = "总结与结论"

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        for item in conclusion:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

    def _add_thank_you_slide(self):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "谢谢观看"

        body_shape = slide.placeholders[1]
        body_shape.text = "如有任何问题，欢迎随时沟通！\n\nQ&A / 联系方式"

    # ==================== 私有方法 - 混合模式 ====================

    def _add_hybrid_title_slide(
        self, title: str, subtitle: Optional[str], background_image: Optional[bytes] = None
    ):
        """混合模式封面：AI 背景 + 文字叠加"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局

        # 背景图
        if background_image:
            try:
                img_stream = io.BytesIO(background_image)
                slide.shapes.add_picture(
                    img_stream,
                    Inches(0),
                    Inches(0),
                    width=self.prs.slide_width,
                    height=self.prs.slide_height,
                )
            except Exception as e:
                logger.warning(f"混合模式背景图失败: {e}")
                # 纯色背景
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor.from_string(self.colors["primary"])
        else:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(self.colors["primary"])

        # 标题文字（半透明背景条）
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.0), Inches(2.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(24)
            p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p2.alignment = PP_ALIGN.CENTER

    def _add_hybrid_content_slide(
        self,
        section: Section,
        background_image: Optional[bytes] = None,
        mermaid_image: Optional[bytes] = None,
    ):
        """混合模式内容页：AI 背景 + 结构化文字"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局

        if background_image:
            try:
                img_stream = io.BytesIO(background_image)
                slide.shapes.add_picture(
                    img_stream,
                    Inches(0),
                    Inches(0),
                    width=self.prs.slide_width,
                    height=self.prs.slide_height,
                )
            except Exception as e:
                logger.warning(f"混合模式背景图失败: {e}")
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor.from_string(self.colors["bg"])
        else:
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(self.colors["bg"])

        # 标题条
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.0), Inches(0.8))
        tf_h = header_box.text_frame
        ph = tf_h.add_paragraph()
        ph.text = section.title
        ph.font.size = Pt(28)
        ph.font.bold = True
        ph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if background_image else RGBColor.from_string(
            self.colors["text"]
        )

        # 文字区域
        has_mermaid = mermaid_image is not None
        text_width = Inches(5.5) if has_mermaid else Inches(11.5)

        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), text_width, Inches(5.5))
        tf = text_box.text_frame
        tf.word_wrap = True

        text_color = (
            RGBColor(0xFF, 0xFF, 0xFF) if background_image else RGBColor.from_string(self.colors["text"])
        )

        for point in section.points:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = text_color
            p.space_after = Pt(6)

        # Mermaid 图表
        if mermaid_image:
            try:
                img_stream = io.BytesIO(mermaid_image)
                slide.shapes.add_picture(img_stream, Inches(6.5), Inches(1.5), width=Inches(6.3))
            except Exception as e:
                logger.warning(f"混合模式 Mermaid 图表失败: {e}")

        # 页码
        slide_idx = list(self.prs.slides).index(slide)
        footer_box = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.4))
        pf = footer_box.text_frame.paragraphs[0]
        pf.text = str(slide_idx + 1)
        pf.font.size = Pt(10)
        pf.font.color.rgb = text_color
        pf.alignment = PP_ALIGN.RIGHT
