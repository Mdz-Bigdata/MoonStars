"""
banana-slides 桥接服务
将 banana-slides-lib 的核心 AI 生成能力适配到现有项目的 PPT 生成链路中。
支持将现有大纲转换为 banana-slides 项目格式，调用 AI 生成全页图片式幻灯片。
"""
import sys
import os
import logging
import base64
from typing import Optional, List, Dict, Any
from pathlib import Path

from app.core.config import settings

logger = logging.getLogger(__name__)

# NOTE: 动态加载 banana-slides-lib 的 Python 模块
BANANA_LIB_PATH = Path(__file__).resolve().parents[4] / "banana-slides-lib" / "backend"


def _ensure_banana_path():
    """确保 banana-slides-lib 的路径在 sys.path 中"""
    lib_str = str(BANANA_LIB_PATH)
    if lib_str not in sys.path:
        sys.path.insert(0, lib_str)


class BananaSlidesService:
    """
    banana-slides 桥接服务

    将 banana-slides 的 AI 服务（Gemini/OpenAI 图片生成）
    和导出服务（python-pptx/PDF）适配到现有 PPT 生成链路。
    """

    def __init__(self):
        self._ai_service = None
        self._export_service = None
        self._available = False
        self._init_error = None
        self._try_init()

    def _try_init(self):
        """尝试初始化 banana-slides 服务"""
        try:
            _ensure_banana_path()

            # 尝试导入 banana-slides 的核心模块
            from services.ai_service import AIService
            from services.export_service import ExportService

            self._ai_service_class = AIService
            self._export_service_class = ExportService
            self._available = True
            logger.info("✅ banana-slides 桥接服务初始化成功")
        except ImportError as e:
            self._init_error = str(e)
            self._available = False
            logger.warning(f"⚠️ banana-slides 桥接不可用: {e}")
        except Exception as e:
            self._init_error = str(e)
            self._available = False
            logger.warning(f"⚠️ banana-slides 初始化异常: {e}")

    @property
    def is_available(self) -> bool:
        """是否可用"""
        return self._available and getattr(settings, "BANANA_SLIDES_ENABLED", False)

    def get_status(self) -> Dict[str, Any]:
        """获取服务状态"""
        return {
            "available": self._available,
            "enabled": getattr(settings, "BANANA_SLIDES_ENABLED", False),
            "error": self._init_error,
            "lib_path": str(BANANA_LIB_PATH),
            "lib_exists": BANANA_LIB_PATH.exists(),
        }

    # ==================== 大纲转换 ====================

    def outline_to_banana_project(
        self,
        title: str,
        subtitle: Optional[str],
        sections: List[Dict[str, Any]],
        theme: str = "light",
    ) -> Dict[str, Any]:
        """
        将现有 PPTOutline 大纲转换为 banana-slides 项目格式。

        Args:
            title: PPT 标题
            subtitle: 副标题
            sections: 章节列表（PPTOutline.sections 序列化后的字典列表）
            theme: 主题名称

        Returns:
            banana-slides 项目格式的字典
        """
        pages = []

        # 封面页
        pages.append({
            "title": title,
            "description": f"封面。标题：{title}。副标题：{subtitle or ''}。",
            "type": "cover",
        })

        # 内容页
        for section in sections:
            points_text = "\n".join(f"- {p}" for p in section.get("points", []))
            desc = f"标题：{section['title']}\n内容要点：\n{points_text}"

            if section.get("mermaid"):
                desc += f"\n\n包含图表：{section.get('visual_type', '图表')}"

            if section.get("table"):
                desc += f"\n\n包含表格数据"

            pages.append({
                "title": section["title"],
                "description": desc,
                "type": "content",
            })

        # 结尾页
        pages.append({
            "title": "谢谢观看",
            "description": "感谢页。Q&A / 联系方式",
            "type": "ending",
        })

        return {
            "title": title,
            "subtitle": subtitle or "",
            "theme": theme,
            "pages": pages,
            "creation_type": "outline",
        }

    # ==================== AI 图片生成 ====================

    async def generate_slide_image(
        self,
        description: str,
        style_prompt: str = "",
        resolution: str = "2k",
    ) -> Optional[bytes]:
        """
        使用 banana-slides 的 AI 服务生成单页幻灯片图片。

        Args:
            description: 页面描述
            style_prompt: 风格提示（可选）
            resolution: 分辨率 '1k' | '2k'

        Returns:
            PNG 图片字节数据，失败返回 None
        """
        if not self.is_available:
            return None

        try:
            _ensure_banana_path()
            from services.ai_service import AIService

            ai_service = AIService()

            # 构建完整的 prompt
            full_prompt = f"""为PPT生成一张精美的幻灯片页面图片。

页面内容描述：
{description}

{f'风格要求：{style_prompt}' if style_prompt else ''}

要求：
- 16:9 宽屏比例
- 专业的排版和设计
- 文字清晰可读
- 配色和谐美观"""

            # 调用 AI 图片生成
            result = ai_service.image_provider.generate_image(
                prompt=full_prompt,
                width=1920 if resolution == "2k" else 1024,
                height=1080 if resolution == "2k" else 576,
            )

            if result and hasattr(result, "content"):
                return result.content
            elif isinstance(result, bytes):
                return result

            return None
        except Exception as e:
            logger.warning(f"banana-slides 图片生成失败: {e}")
            return None

    # ==================== 批量生成幻灯片 ====================

    async def generate_all_slides(
        self,
        project_data: Dict[str, Any],
        progress_callback=None,
    ) -> List[Optional[bytes]]:
        """
        批量生成所有幻灯片图片。

        Args:
            project_data: banana-slides 项目格式数据
            progress_callback: 进度回调 (current, total, message)

        Returns:
            图片字节数据列表（失败的页面为 None）
        """
        pages = project_data.get("pages", [])
        results = []

        for i, page in enumerate(pages):
            if progress_callback:
                await progress_callback(i, len(pages), f"正在 AI 生成第 {i + 1}/{len(pages)} 页...")

            image = await self.generate_slide_image(
                description=page.get("description", page.get("title", "")),
                style_prompt=project_data.get("theme", ""),
            )
            results.append(image)

        return results

    # ==================== 导出服务 ====================

    def export_pptx_from_images(self, image_data_list: List[bytes]) -> Optional[bytes]:
        """
        使用 banana-slides 的导出服务将图片列表组装为 PPTX。

        Args:
            image_data_list: 图片字节数据列表

        Returns:
            PPTX 文件字节数据
        """
        if not self._available:
            return None

        try:
            import tempfile
            import io
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            for img_data in image_data_list:
                if img_data is None:
                    continue

                slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
                img_stream = io.BytesIO(img_data)
                slide.shapes.add_picture(
                    img_stream,
                    Inches(0),
                    Inches(0),
                    width=prs.slide_width,
                    height=prs.slide_height,
                )

            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            return output.read()
        except Exception as e:
            logger.error(f"PPTX 导出失败: {e}")
            return None
