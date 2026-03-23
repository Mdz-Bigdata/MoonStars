"""
文档处理服务模块
整合 PDF/PPT/Office 文档转换功能

功能:
- PDF 转 Markdown（使用 MinerU CLI）
- PPT 转网页 HTML
- Office 文档转 Markdown
- URL 转 PPT
"""
import os
import sys
import shutil
import logging
import tempfile
import subprocess
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional, Dict, Any, List

# PDF 处理 - 优先使用 pdfplumber
try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

# PDF 处理 - PyMuPDF 备选
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

# PPT 处理
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# 图片转 PDF
try:
    import img2pdf
    HAS_IMG2PDF = True
except ImportError:
    HAS_IMG2PDF = False

from PIL import Image

from app.core.config import settings

logger = logging.getLogger(__name__)


@dataclass
class DocumentProcessorConfig:
    """文档处理配置"""
    # MinerU 配置
    backend: str = "hybrid-auto-engine"
    language: str = "zh"  # 默认中文，避免乱码
    dpi: int = 150
    quality: int = 85
    
    # PPT 配置
    slide_width_inches: float = 10.0
    slide_height_inches: float = 7.5
    
    # 临时目录
    temp_dir: str = ""
    
    def __post_init__(self):
        if not self.temp_dir:
            self.temp_dir = os.path.join(settings.UPLOAD_DIR, "temp", "doc_processor")
        os.makedirs(self.temp_dir, exist_ok=True)


class DocumentProcessor:
    """文档处理服务"""
    
    def __init__(self, config: Optional[DocumentProcessorConfig] = None):
        self.config = config or DocumentProcessorConfig()
        self._temp_files: List[str] = []
    
    async def process_pdf_to_markdown(self, file_path: str) -> str:
        """
        将 PDF 转换为 Markdown
        优先使用 MinerU CLI，不可用时使用高保真图片渲染
        """
        logger.info(f"📄 开始 PDF 转换: {file_path}")
        
        output_dir = os.path.join(self.config.temp_dir, f"pdf_{os.getpid()}")
        os.makedirs(output_dir, exist_ok=True)
        self._temp_files.append(output_dir)
        
        # 检测可用的 MinerU 命令
        mineru_cmd = self._find_mineru_command()
        if mineru_cmd:
            # 使用 MinerU CLI
            result = await self._process_pdf_with_mineru(file_path, output_dir, mineru_cmd)
            if result and not result.startswith("转换失败"):
                return result
            logger.warning("MinerU 处理失败，尝试高保真图片渲染")
        else:
            logger.warning("MinerU CLI 未找到，使用高保真图片渲染模式")
        
        # 高保真回退方案：将每页渲染为图片
        return await self._process_pdf_as_images(file_path, output_dir)
    
    async def _process_pdf_with_mineru(self, file_path: str, output_dir: str, mineru_cmd: str) -> str:
        """使用 MinerU CLI 处理 PDF"""
        abs_file_path = os.path.abspath(file_path)
        abs_output_dir = os.path.abspath(output_dir)
        
        cmd = [
            mineru_cmd,
            "pdf",  # 新版本使用子命令格式
            "--pdf", abs_file_path,
            "--method", "auto"
        ]
        
        try:
            env = os.environ.copy()
            env["PYTHONIOENCODING"] = "utf-8"
            env["KMP_DUPLICATE_LIB_OK"] = "TRUE"
            
            # 使用 MinerU 时显式指定输出目录
            # 某些版本的 magic-pdf 必须使用 --output-dir 来指定输出路径
            # 如果不指定，它可能会默认输出到 /tmp/magic-pdf 导致我们找不到文件
            cmd = [
                mineru_cmd,
                "pdf",
                "--pdf", abs_file_path,
                "--method", "auto",
                "--output-dir", abs_output_dir
            ]
            
            logger.info(f"执行命令: {' '.join(cmd)}")
            
            result = subprocess.run(
                cmd,
                capture_output=True,
                timeout=900,
                env=env,
                cwd=abs_output_dir  # 将工作目录设置到临时输出目录，某些版本会在 cwd 生成
            )
            
            stdout = self._safe_decode(result.stdout)
            stderr = self._safe_decode(result.stderr)
            
            if stdout:
                logger.info(f"MinerU stdout: {stdout}")
            if stderr:
                logger.info(f"MinerU stderr: {stderr}")
            
            if result.returncode != 0:
                logger.error(f"MinerU 执行失败，返回码: {result.returncode}")
                # 将完整的错误信息包含在异常中，便于排查
                error_msg = stderr if stderr else "未知错误"
                return f"转换失败：MinerU 错误 (code={result.returncode}) - {error_msg[:200]}"
            
            # 搜索 MinerU 生成的 Markdown
            # 新版 MinerU 通常会在指定的 output-dir 下新建一个同名去后缀的文件夹
            pdf_dir = os.path.dirname(abs_file_path)
            pdf_stem = Path(abs_file_path).stem
            
            search_dirs = [
                abs_output_dir, # 当前执行目录
                os.path.join(abs_output_dir, pdf_stem), # 执行目录下的同名文件夹
                pdf_dir,        # PDF 所在目录
                os.path.join(pdf_dir, pdf_stem), # PDF 同名文件夹
                f"/tmp/magic-pdf/{pdf_stem}/auto", # 兜底寻找默认 tmp 目录
            ]
            
            markdown_content = None
            for s_dir in search_dirs:
                if os.path.exists(s_dir):
                    markdown_content = self._find_and_read_markdown(s_dir)
                    if markdown_content:
                        break
            
            if not markdown_content:
                logger.error("未找到生成的 Markdown 文件")
                return "转换失败：未生成 Markdown 文件"
            
            logger.info(f"✅ MinerU PDF 转换成功，内容长度: {len(markdown_content)}")
            return markdown_content
            
        except subprocess.TimeoutExpired:
            logger.error("MinerU 执行超时")
            return "转换失败：处理超时"
        except Exception as e:
            logger.error(f"MinerU 处理异常: {e}")
            return f"转换失败：{str(e)}"
    
    async def _process_pdf_with_pymupdf(self, file_path: str) -> str:
        """使用 pdfplumber 或 PyMuPDF 提取 PDF 文本（回退方案）"""
        # 优先使用 pdfplumber
        if HAS_PDFPLUMBER:
            return await self._process_pdf_with_pdfplumber(file_path)
        
        # 备选使用 PyMuPDF
        if HAS_PYMUPDF:
            return await self._process_pdf_with_fitz(file_path)
        
        logger.error("PDF 处理库未安装（需要 pdfplumber 或 PyMuPDF）")
        return "转换失败：PDF 处理库未安装"
    
    async def _process_pdf_with_pdfplumber(self, file_path: str) -> str:
        """使用 pdfplumber 提取 PDF 文本"""
        try:
            markdown_parts = []
            title = Path(file_path).stem
            markdown_parts.append(f"# {title}\n")
            
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    
                    if text.strip():
                        markdown_parts.append(f"\n## 第 {page_num + 1} 页\n")
                        paragraphs = text.split('\n\n')
                        for para in paragraphs:
                            cleaned = para.strip()
                            if cleaned:
                                markdown_parts.append(cleaned + "\n")
            
            content = "\n".join(markdown_parts)
            logger.info(f"✅ pdfplumber PDF 转换成功，内容长度: {len(content)}")
            return content
            
        except Exception as e:
            logger.error(f"pdfplumber 处理异常: {e}")
            return f"转换失败：{str(e)}"
    
    async def _process_pdf_with_fitz(self, file_path: str) -> str:
        """使用 PyMuPDF (fitz) 提取 PDF 文本"""
        try:
            doc = fitz.open(file_path)
            markdown_parts = []
            
            title = Path(file_path).stem
            markdown_parts.append(f"# {title}\n")
            
            for page_num, page in enumerate(doc):
                text = page.get_text("text")
                
                if text.strip():
                    markdown_parts.append(f"\n## 第 {page_num + 1} 页\n")
                    paragraphs = text.split('\n\n')
                    for para in paragraphs:
                        cleaned = para.strip()
                        if cleaned:
                            markdown_parts.append(cleaned + "\n")
            
            doc.close()
            
            content = "\n".join(markdown_parts)
            logger.info(f"✅ PyMuPDF PDF 转换成功，内容长度: {len(content)}")
            return content
            
        except Exception as e:
            logger.error(f"PyMuPDF 处理异常: {e}")
            return f"转换失败：{str(e)}"
    
    async def _process_pdf_as_images(self, file_path: str, output_dir: str) -> str:
        """
        高保真模式：将 PDF 每页渲染为图片
        使用 pdftoppm 命令行工具（Poppler）渲染高质量图片
        """
        try:
            abs_file_path = os.path.abspath(file_path)
            title = Path(file_path).stem
            
            # 创建图片输出目录
            images_dir = os.path.join(output_dir, "images")
            os.makedirs(images_dir, exist_ok=True)
            
            # 使用 pdftoppm 渲染 PDF 为高分辨率 PNG 图片
            pdftoppm_path = shutil.which("pdftoppm")
            if not pdftoppm_path:
                logger.error("pdftoppm 未找到，请安装 Poppler")
                # 回退到文本提取
                return await self._process_pdf_with_pymupdf(file_path)
            
            # 渲染命令：200 DPI，PNG 格式
            cmd = [
                pdftoppm_path,
                "-png",
                "-r", str(self.config.dpi),
                abs_file_path,
                os.path.join(images_dir, "page")
            ]
            
            logger.info(f"执行 pdftoppm: {' '.join(cmd)}")
            
            result = subprocess.run(
                cmd,
                capture_output=True,
                timeout=300
            )
            
            if result.returncode != 0:
                stderr = self._safe_decode(result.stderr)
                logger.error(f"pdftoppm 失败: {stderr}")
                return await self._process_pdf_with_pymupdf(file_path)
            
            # 查找生成的图片文件
            image_files = sorted([
                f for f in os.listdir(images_dir) 
                if f.endswith('.png')
            ])
            
            if not image_files:
                logger.error("未生成任何图片")
                return await self._process_pdf_with_pymupdf(file_path)
            
            # 将图片复制到 uploads 目录并生成 Markdown
            markdown_parts = [f"# {title}\n"]
            
            for idx, img_file in enumerate(image_files):
                src_path = os.path.join(images_dir, img_file)
                
                # 复制到 uploads/pdf_images 目录
                pdf_images_dir = os.path.join(settings.UPLOAD_DIR, "pdf_images")
                os.makedirs(pdf_images_dir, exist_ok=True)
                
                # 使用唯一文件名
                import time
                unique_name = f"{title}_{int(time.time())}_{idx + 1}.png"
                dst_path = os.path.join(pdf_images_dir, unique_name)
                shutil.copy2(src_path, dst_path)
                
                # 生成相对 URL
                img_url = f"/uploads/pdf_images/{unique_name}"
                
                # 添加到 Markdown，每页作为一个图片
                markdown_parts.append(f"\n![第 {idx + 1} 页]({img_url})\n")
            
            content = "\n".join(markdown_parts)
            logger.info(f"✅ 高保真 PDF 转换成功，共 {len(image_files)} 页")
            return content
            
        except subprocess.TimeoutExpired:
            logger.error("pdftoppm 执行超时")
            return await self._process_pdf_with_pymupdf(file_path)
        except Exception as e:
            logger.error(f"高保真 PDF 处理异常: {e}")
            return await self._process_pdf_with_pymupdf(file_path)
    
    async def process_office_doc(self, file_path: str) -> str:
        """
        处理 Office 文档（Word, Excel）
        先转换为 PDF，再用 MinerU 处理
        """
        logger.info(f"📑 开始 Office 文档处理: {file_path}")
        
        ext = Path(file_path).suffix.lower()
        
        # 尝试使用 LibreOffice 转换为 PDF
        pdf_path = await self._convert_to_pdf_with_libreoffice(file_path)
        
        if pdf_path and os.path.exists(pdf_path):
            # 用 MinerU 处理转换后的 PDF
            return await self.process_pdf_to_markdown(pdf_path)
        
        # 回退方案：直接读取文本
        if ext in ['.docx', '.doc']:
            return await self._extract_docx_text(file_path)
        elif ext in ['.xlsx', '.xls']:
            return await self._extract_xlsx_text(file_path)
        
        return "转换失败：不支持的文档格式"
    
    async def ppt_to_web(self, file_path: str) -> str:
        """
        将 PPT/PPTX 转换为网页 HTML（高保真模式）
        使用 Playwright 渲染每页幻灯片为图片
        """
        logger.info(f"🎯 开始 PPT 高保真转换: {file_path}")
        
        output_dir = os.path.join(self.config.temp_dir, f"ppt_{os.getpid()}")
        os.makedirs(output_dir, exist_ok=True)
        self._temp_files.append(output_dir)
        
        # 方案一：使用 LibreOffice/soffice 转 PDF 再渲染为图片
        pdf_path = await self._convert_to_pdf_with_libreoffice(file_path)
        
        if pdf_path and os.path.exists(pdf_path):
            logger.info(f"PPT 已转换为 PDF: {pdf_path}")
            return await self._process_pdf_as_images(pdf_path, output_dir)
        
        # 方案二：使用 Playwright 渲染为图片
        logger.warning("LibreOffice 不可用，尝试 Playwright 渲染")
        result = await self._render_ppt_with_playwright(file_path, output_dir)
        if result:
            return result
        
        # 方案三：回退到文本提取
        logger.warning("Playwright 渲染失败，回退到文本提取")
        return await self._extract_ppt_text(file_path, output_dir)
    
    async def _render_ppt_with_playwright(self, file_path: str, output_dir: str) -> Optional[str]:
        """使用 Playwright 渲染 PPT 每页为图片"""
        if not HAS_PPTX:
            logger.error("python-pptx 未安装")
            return None
        
        try:
            from playwright.async_api import async_playwright
            from pptx.util import Inches, Pt, Emu
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            import base64
            from io import BytesIO
            
            prs = Presentation(file_path)
            slide_width = prs.slide_width.emu / 914400  # 转换为英寸
            slide_height = prs.slide_height.emu / 914400
            
            # 计算像素尺寸（96 DPI）
            px_width = int(slide_width * 96 * 2)  # 2x 高清
            px_height = int(slide_height * 96 * 2)
            
            images_dir = os.path.join(output_dir, "images")
            os.makedirs(images_dir, exist_ok=True)
            
            title = Path(file_path).stem
            markdown_parts = [f"# {title}\n"]
            
            async with async_playwright() as p:
                browser = await p.chromium.launch(headless=True)
                
                for idx, slide in enumerate(prs.slides):
                    try:
                        # 生成幻灯片 HTML
                        html_content = self._slide_to_html(slide, prs, idx + 1)
                        
                        # 使用 Playwright 渲染并截图
                        page = await browser.new_page(viewport={"width": px_width, "height": px_height})
                        await page.set_content(html_content)
                        await page.wait_for_timeout(500)  # 等待渲染
                        
                        # 截图保存
                        screenshot_path = os.path.join(images_dir, f"slide_{idx + 1:03d}.png")
                        await page.screenshot(path=screenshot_path, full_page=False)
                        await page.close()
                        
                        # 移动到上传目录
                        static_dir = os.path.join("./uploads", "ppt_slides")
                        os.makedirs(static_dir, exist_ok=True)
                        final_path = os.path.join(static_dir, f"{title}_slide_{idx + 1:03d}.png")
                        shutil.copy2(screenshot_path, final_path)
                        
                        # 添加到 Markdown
                        image_url = f"/uploads/ppt_slides/{title}_slide_{idx + 1:03d}.png"
                        markdown_parts.append(f"\n## 第 {idx + 1} 页\n\n![第 {idx + 1} 页]({image_url})\n")
                        
                        logger.info(f"✅ 幻灯片 {idx + 1} 截图成功")
                        
                    except Exception as e:
                        logger.error(f"幻灯片 {idx + 1} 渲染失败: {e}")
                        # 回退到文本
                        texts = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                texts.append(shape.text.strip())
                        markdown_parts.append(f"\n## 第 {idx + 1} 页\n\n" + "\n\n".join(texts) + "\n")
                
                await browser.close()
            
            content = "\n".join(markdown_parts)
            logger.info(f"✅ PPT Playwright 渲染成功，共 {len(prs.slides)} 页")
            return content
            
        except Exception as e:
            logger.error(f"Playwright PPT 渲染异常: {e}")
            import traceback
            logger.error(traceback.format_exc())
            return None
    
    def _slide_to_html(self, slide, prs, slide_num: int) -> str:
        """将单个幻灯片转换为 HTML - 增强版"""
        from pptx.util import Emu
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.oxml.ns import qn
        import base64
        from io import BytesIO
        import zipfile
        
        slide_width = prs.slide_width.emu / 914400 * 96  # 转为像素
        slide_height = prs.slide_height.emu / 914400 * 96
        
        # 提取背景图片或背景样式
        bg_style = self._extract_slide_background(slide, prs)
        
        html_parts = [f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            width: {int(slide_width * 2)}px;
            height: {int(slide_height * 2)}px;
            {bg_style}
            font-family: "Microsoft YaHei", "SimHei", "PingFang SC", Arial, sans-serif;
            position: relative;
            overflow: hidden;
        }}
        .shape {{
            position: absolute;
            display: flex;
            flex-direction: column;
            white-space: pre-wrap;
            word-wrap: break-word;
            overflow: hidden;
        }}
        .shape img {{
            width: 100%;
            height: 100%;
            object-fit: contain;
        }}
        p {{ margin: 0; padding: 0; }}
    </style>
</head>
<body>
''']
        
        # 递归处理所有形状（包括组合形状）
        self._render_shapes(slide.shapes, html_parts, 0, 0)
        
        html_parts.append('</body></html>')
        return "".join(html_parts)
    
    def _extract_slide_background(self, slide, prs) -> str:
        """提取幻灯片背景样式"""
        import base64
        from pptx.enum.dml import MSO_THEME_COLOR
        from pptx.oxml.ns import qn
        
        try:
            # 检查幻灯片是否有独立背景
            bg = slide.background
            fill = bg.fill
            
            # 检查背景填充类型
            if fill.type is not None:
                # 纯色背景
                if fill.type == 1:  # MSO_FILL.SOLID
                    try:
                        if fill.fore_color and fill.fore_color.rgb:
                            return f"background-color: #{fill.fore_color.rgb};"
                    except:
                        pass
                
                # 图片背景
                if fill.type == 6:  # MSO_FILL.PICTURE
                    try:
                        # 尝试从 background 的 blipFill 获取图片
                        bg_elem = slide._element.cSld.bg
                        if bg_elem is not None:
                            blip_fill = bg_elem.find('.//a:blipFill', {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                            if blip_fill is not None:
                                blip = blip_fill.find('.//a:blip', {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                                if blip is not None:
                                    embed_id = blip.get(qn('r:embed'))
                                    if embed_id:
                                        rel = slide.part.related_parts.get(embed_id)
                                        if rel:
                                            img_bytes = rel.blob
                                            img_type = rel.content_type.split('/')[-1]
                                            img_b64 = base64.b64encode(img_bytes).decode()
                                            return f"background-image: url('data:image/{img_type};base64,{img_b64}'); background-size: cover; background-position: center;"
                    except Exception as e:
                        logger.debug(f"背景图片提取失败: {e}")
                
                # 渐变背景
                if fill.type == 2:  # MSO_FILL.GRADIENT
                    try:
                        # 尝试获取渐变颜色
                        colors = []
                        if hasattr(fill, 'gradient_stops'):
                            for stop in fill.gradient_stops:
                                if stop.color and stop.color.rgb:
                                    colors.append(f"#{stop.color.rgb}")
                        if len(colors) >= 2:
                            return f"background: linear-gradient(135deg, {', '.join(colors)});"
                    except:
                        pass
            
            # 检查幻灯片布局的背景
            try:
                layout_bg = slide.slide_layout.background.fill
                if layout_bg.type == 6:  # 图片
                    # 尝试从布局获取背景
                    pass
            except:
                pass
                
        except Exception as e:
            logger.debug(f"背景提取异常: {e}")
        
        return "background-color: #FFFFFF;"
    
    def _render_shapes(self, shapes, html_parts, offset_left, offset_top):
        """递归渲染形状（支持组合形状）"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        import base64
        
        for shape in shapes:
            try:
                # 计算位置和大小
                left = (shape.left / 914400 * 96 * 2 if shape.left else 0) + offset_left
                top = (shape.top / 914400 * 96 * 2 if shape.top else 0) + offset_top
                width = shape.width / 914400 * 96 * 2 if shape.width else 100
                height = shape.height / 914400 * 96 * 2 if shape.height else 50
                
                # 处理组合形状
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    try:
                        self._render_shapes(shape.shapes, html_parts, left, top)
                    except:
                        pass
                    continue
                
                # 处理图片
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        img_bytes = image.blob
                        img_ext = image.content_type.split('/')[-1]
                        img_b64 = base64.b64encode(img_bytes).decode()
                        html_parts.append(f'''
    <div class="shape" style="left:{left:.0f}px; top:{top:.0f}px; width:{width:.0f}px; height:{height:.0f}px;">
        <img src="data:image/{img_ext};base64,{img_b64}" />
    </div>
''')
                    except Exception as e:
                        logger.debug(f"图片处理失败: {e}")
                    continue
                
                # 处理占位符中的图片
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    try:
                        if hasattr(shape, 'image') and shape.image:
                            image = shape.image
                            img_bytes = image.blob
                            img_ext = image.content_type.split('/')[-1]
                            img_b64 = base64.b64encode(img_bytes).decode()
                            html_parts.append(f'''
    <div class="shape" style="left:{left:.0f}px; top:{top:.0f}px; width:{width:.0f}px; height:{height:.0f}px;">
        <img src="data:image/{img_ext};base64,{img_b64}" />
    </div>
''')
                            continue
                    except:
                        pass
                        
                # 处理文本框
                if hasattr(shape, "text_frame"):
                    text_frame = shape.text_frame
                    if text_frame and text_frame.paragraphs:
                        text_html = []
                        has_content = False
                        
                        for para in text_frame.paragraphs:
                            para_html = []
                            for run in para.runs:
                                run_text = run.text or ""
                                if not run_text:
                                    continue
                                has_content = True
                                
                                # 提取样式
                                styles = []
                                if run.font.size:
                                    styles.append(f"font-size:{run.font.size.pt * 2:.0f}px")
                                else:
                                    styles.append("font-size:24px")
                                    
                                if run.font.color and run.font.color.rgb:
                                    styles.append(f"color:#{run.font.color.rgb}")
                                else:
                                    styles.append("color:#000000")
                                    
                                if run.font.bold:
                                    styles.append("font-weight:bold")
                                if run.font.italic:
                                    styles.append("font-style:italic")
                                    
                                style_str = "; ".join(styles)
                                para_html.append(f'<span style="{style_str}">{run_text}</span>')
                            
                            if para_html:
                                text_html.append(f"<p>{''.join(para_html)}</p>")
                        
                        if has_content and text_html:
                            html_parts.append(f'''
    <div class="shape" style="left:{left:.0f}px; top:{top:.0f}px; width:{width:.0f}px; height:{height:.0f}px; justify-content:flex-start; padding:8px;">
        {"".join(text_html)}
    </div>
''')
                    
            except Exception as e:
                logger.debug(f"形状渲染失败: {e}")
                continue
    
    async def _extract_ppt_text(self, file_path: str, output_dir: str) -> str:
        """从 PPT 提取文本（回退方案）"""
        if not HAS_PPTX:
            return "转换失败：python-pptx 未安装"
        
        try:
            prs = Presentation(file_path)
            slides_html = []
            
            for idx, slide in enumerate(prs.slides):
                slide_content = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                slides_html.append(f"## 第 {idx + 1} 页\n\n" + "\n\n".join(slide_content))
            
            markdown = f"# {Path(file_path).stem}\n\n" + "\n\n---\n\n".join(slides_html)
            logger.info(f"✅ PPT 文本提取成功，共 {len(prs.slides)} 页")
            return markdown
            
        except Exception as e:
            logger.error(f"PPT 处理异常: {e}")
            return f"转换失败：{str(e)}"
    
    async def url_to_ppt(self, url: str, config: Optional[Dict[str, Any]] = None) -> Optional[str]:
        """
        将网页 URL 转换为 PPT
        通过截图网页内容并转换为 PPT 格式
        """
        logger.info(f"🌐 开始 URL 转 PPT: {url}")
        
        if not HAS_PPTX:
            logger.error("python-pptx 未安装")
            return None
        
        if not HAS_PYMUPDF:
            logger.error("PyMuPDF 未安装")
            return None
        
        work_dir = os.path.join(self.config.temp_dir, f"url_ppt_{os.getpid()}")
        os.makedirs(work_dir, exist_ok=True)
        self._temp_files.append(work_dir)
        
        try:
            # 使用 playwright 截取网页
            screenshots = await self._capture_webpage_screenshots(url, work_dir)
            
            if not screenshots:
                logger.error("网页截图失败")
                return None
            
            # 创建 PPT
            prs = Presentation()
            prs.slide_width = Inches(self.config.slide_width_inches)
            prs.slide_height = Inches(self.config.slide_height_inches)
            
            blank_layout = prs.slide_layouts[6]  # 空白布局
            
            for img_path in screenshots:
                slide = prs.slides.add_slide(blank_layout)
                
                # 添加图片，铺满整个幻灯片
                slide.shapes.add_picture(
                    img_path,
                    Inches(0),
                    Inches(0),
                    prs.slide_width,
                    prs.slide_height
                )
            
            # 保存 PPT
            import time
            output_filename = f"web_ppt_{int(time.time())}.pptx"
            output_path = os.path.join(settings.UPLOAD_DIR, output_filename)
            prs.save(output_path)
            
            logger.info(f"✅ URL 转 PPT 成功: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"URL 转 PPT 异常: {e}")
            return None
    
    async def process_text_file(self, file_path: str) -> str:
        """
        处理文本文件（TXT, MD, CSV, JSON）
        """
        logger.info(f"📝 处理文本文件: {file_path}")
        
        try:
            # 尝试多种编码读取
            for encoding in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    logger.info(f"使用 {encoding} 编码读取成功")
                    return content
                except UnicodeDecodeError:
                    continue
            
            # 最后尝试二进制读取
            with open(file_path, 'rb') as f:
                content = f.read().decode('utf-8', errors='replace')
            return content
            
        except Exception as e:
            logger.error(f"文本文件处理异常: {e}")
            return f"转换失败：{str(e)}"
    
    def cleanup_temp(self):
        """清理临时文件"""
        for path in self._temp_files:
            try:
                if os.path.isdir(path):
                    shutil.rmtree(path, ignore_errors=True)
                elif os.path.isfile(path):
                    os.remove(path)
            except Exception as e:
                logger.warning(f"清理临时文件失败: {path}, {e}")
        
        self._temp_files.clear()
    
    # ==================== 私有方法 ====================
    
    def _find_mineru_command(self) -> Optional[str]:
        """查找可用的 MinerU 命令"""
        possible_paths = [
            # 项目虚拟环境
            os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "venv/bin"),
            # 用户 MinerU 环境
            "/Users/mindezhi/mineru_env/bin",
            "/Users/mindezhi/miniconda3/bin",
            os.path.expanduser("~/mineru_env/bin"),
            # 系统路径
            "/opt/homebrew/bin",
            "/usr/local/bin"
        ]
        
        # 构建搜索路径
        search_path = ":".join(possible_paths) + ":" + os.environ.get("PATH", "")
        
        # 尝试查找 mineru 或 magic-pdf
        for cmd in ["mineru", "magic-pdf"]:
            cmd_path = shutil.which(cmd, path=search_path)
            if cmd_path:
                logger.info(f"找到 MinerU 命令: {cmd_path}")
                return cmd_path
        
        return None
    
    def _safe_decode(self, data: bytes) -> str:
        """安全解码字节数据为 UTF-8 字符串"""
        if not data:
            return ""
        try:
            return data.decode('utf-8')
        except UnicodeDecodeError:
            try:
                return data.decode('gbk')
            except:
                return data.decode('latin-1', errors='replace')
    
    def _find_and_read_markdown(self, output_dir: str) -> Optional[str]:
        """在输出目录中查找并读取 Markdown 文件"""
        md_files = []
        
        for root, dirs, files in os.walk(output_dir):
            for f in files:
                if f.endswith('.md'):
                    md_files.append(os.path.join(root, f))
        
        if not md_files:
            return None
        
        # 读取第一个 Markdown 文件
        md_path = md_files[0]
        logger.info(f"读取 Markdown 文件: {md_path}")
        
        try:
            with open(md_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logger.error(f"读取 Markdown 失败: {e}")
            return None
    
    async def _convert_to_pdf_with_libreoffice(self, file_path: str) -> Optional[str]:
        """使用 LibreOffice/soffice 转换文档为 PDF"""
        output_dir = os.path.join(self.config.temp_dir, "libreoffice_output")
        os.makedirs(output_dir, exist_ok=True)
        
        # 查找可用的 LibreOffice 命令 - 按优先级排序
        possible_paths = [
            # 标准安装位置
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            # 用户目录安装
            os.path.expanduser("~/Applications/LibreOffice.app/Contents/MacOS/soffice"),
            # DMG 挂载位置
            "/Volumes/LibreOffice/LibreOffice.app/Contents/MacOS/soffice",
            # Homebrew 链接
            "/opt/homebrew/bin/soffice",
            "/usr/local/bin/soffice",
            # Linux 路径
            "/usr/bin/soffice",
            "/usr/bin/libreoffice",
        ]
        
        soffice_cmd = None
        for path in possible_paths:
            if os.path.exists(path) and os.access(path, os.X_OK):
                soffice_cmd = path
                logger.info(f"找到 LibreOffice: {path}")
                break
        
        # 如果固定路径未找到，尝试 which
        if not soffice_cmd:
            for cmd_name in ["soffice", "libreoffice"]:
                cmd_path = shutil.which(cmd_name)
                if cmd_path:
                    soffice_cmd = cmd_path
                    break
        
        if not soffice_cmd:
            logger.warning("LibreOffice/soffice 未找到，已检查的路径: " + ", ".join(possible_paths[:4]))
            return None
        
        try:
            abs_file_path = os.path.abspath(file_path)
            abs_output_dir = os.path.abspath(output_dir)
            
            cmd = [
                soffice_cmd,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", abs_output_dir,
                abs_file_path
            ]
            
            logger.info(f"执行 LibreOffice: {' '.join(cmd)}")
            
            result = subprocess.run(cmd, capture_output=True, timeout=180)
            
            stdout = self._safe_decode(result.stdout)
            stderr = self._safe_decode(result.stderr)
            
            if stdout:
                logger.info(f"LibreOffice stdout: {stdout[:300]}")
            if stderr:
                logger.warning(f"LibreOffice stderr: {stderr[:300]}")
            
            if result.returncode == 0:
                pdf_name = Path(file_path).stem + ".pdf"
                pdf_path = os.path.join(abs_output_dir, pdf_name)
                if os.path.exists(pdf_path):
                    logger.info(f"✅ LibreOffice 转换成功: {pdf_path}")
                    return pdf_path
                else:
                    logger.warning(f"PDF 文件未生成: {pdf_path}")
            else:
                logger.warning(f"LibreOffice 返回码: {result.returncode}")
        except subprocess.TimeoutExpired:
            logger.warning("LibreOffice 转换超时")
        except Exception as e:
            logger.warning(f"LibreOffice 转换失败: {e}")
        
        return None
    
    async def _extract_docx_text(self, file_path: str) -> str:
        """从 DOCX 提取文本"""
        try:
            from docx import Document
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            logger.error(f"DOCX 提取失败: {e}")
            return f"转换失败：{str(e)}"
    
    async def _extract_xlsx_text(self, file_path: str) -> str:
        """从 XLSX 提取文本"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            content = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                content.append(f"## {sheet_name}\n")
                
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_values):
                        rows.append("| " + " | ".join(row_values) + " |")
                
                if rows:
                    # 添加表头分隔符
                    header_sep = "| " + " | ".join(["---"] * len(rows[0].split("|")[1:-1])) + " |"
                    content.append(rows[0])
                    content.append(header_sep)
                    content.extend(rows[1:])
                
                content.append("")
            
            return "\n".join(content)
        except Exception as e:
            logger.error(f"XLSX 提取失败: {e}")
            return f"转换失败：{str(e)}"
    
    async def _capture_webpage_screenshots(self, url: str, output_dir: str) -> List[str]:
        """使用 Playwright 截取网页截图"""
        try:
            from playwright.async_api import async_playwright
            
            screenshots = []
            
            async with async_playwright() as p:
                browser = await p.chromium.launch(headless=True)
                page = await browser.new_page(viewport={"width": 1920, "height": 1080})
                
                await page.goto(url, wait_until="networkidle", timeout=60000)
                await page.wait_for_timeout(3000)  # 等待渲染
                
                # 获取页面高度
                page_height = await page.evaluate("document.body.scrollHeight")
                viewport_height = 1080
                
                # 分页截图
                scroll_position = 0
                page_num = 0
                
                while scroll_position < page_height:
                    await page.evaluate(f"window.scrollTo(0, {scroll_position})")
                    await page.wait_for_timeout(500)
                    
                    screenshot_path = os.path.join(output_dir, f"page_{page_num:03d}.png")
                    await page.screenshot(path=screenshot_path, full_page=False)
                    screenshots.append(screenshot_path)
                    
                    scroll_position += viewport_height
                    page_num += 1
                    
                    # 防止无限循环
                    if page_num > 50:
                        break
                
                await browser.close()
            
            return screenshots
            
        except ImportError:
            logger.error("Playwright 未安装")
            return []
        except Exception as e:
            logger.error(f"网页截图失败: {e}")
            return []


# 创建全局实例
doc_processor = DocumentProcessor()
